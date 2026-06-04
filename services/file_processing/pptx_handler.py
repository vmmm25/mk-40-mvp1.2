import logging
from pathlib import Path
from .utils import get_output_path, gemini_client

logger = logging.getLogger(__name__)

def process_pptx(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "summarize"

    def _read_pptx_text() -> str:
        try:
            from pptx import Presentation
            prs  = Presentation(path)
            text = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = f"\n--- Slide {i} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text.strip() + "\n"
                text.append(slide_text)
            return "\n".join(text)
        except ImportError:
            return "python-pptx not installed."

    if action in ("summarize", "extract_text", "analyze"):
        text = _read_pptx_text()
        if action == "extract_text":
            out = get_output_path(path, "text", ".txt")
            out.write_text(text, encoding="utf-8")
            return f"Text extracted. Saved: {out.name}"
        try:
            model    = gemini_client()
            prompt   = f"{'Summarize' if action == 'summarize' else 'Analyze'} this presentation:\n{text[:30000]}"
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            logger.exception("AI processing failed")
            return f"AI processing failed: {e}"

    if action == "to_pdf":
        return "Convert to PDF requires LibreOffice/PowerPoint integration."

    return f"Unknown PPTX action: '{action}'. Try: summarize, extract_text, analyze"
